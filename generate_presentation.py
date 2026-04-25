
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    def add_slide(title, points, image_path=None):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            if ">" in point: # Sub-points
                p.text = point.replace(">", "").strip()
                p.level = 1

        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), height=Inches(3.5))

    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Hybrid AI Intelligence Dashboard"
    subtitle.text = "Detecting Framing Bias, Fake News, and Manipulation\n\nSUBMITTED BY: GROUP NO: AI-07\nSymbiosis Skills & Professional University, Pune"

    # 2. CONTENTS
    add_slide("CONTENTS", [
        "INTRODUCTION", "LITERATURE REVIEW", "SUBJECT INTEGRATION OVERVIEW",
        "DESIGN AND METHODOLOGY", "SYSTEM ARCHITECTURE", "IMPLEMENTATION",
        "RESULTS", "INDUSTRY RELEVANCE", "CHALLENGES",
        "ETHICAL CONSIDERATIONS", "FUTURE SCOPE", "CONCLUSION", "REFERENCES"
    ])

    # 3. INTRODUCTION
    add_slide("INTRODUCTION", [
        "In the era of post-truth politics, information is weaponized.",
        "Framing Bias: Selective presentation of facts to steer narratives.",
        "Problem: Traditional filters fail to detect subtle linguistic manipulation.",
        "Solution: A Hybrid AI approach combining local ML with Cloud LLMs.",
        "Goal: Real-time transparency and cognitive defense."
    ])

    # 4. LITERATURE REVIEW
    add_slide("LITERATURE REVIEW", [
        "Manual Fact-Checking: High accuracy but high latency (Snopes, PolitiFact).",
        "Traditional ML: Fast but lacks semantic depth (TF-IDF + Naive Bayes).",
        "State-of-the-Art: LLMs provide reasoning but lack live context.",
        "Our Innovation: 'Dynamic Live Context' integration for near-real-time verification."
    ])

    # 5. SUBJECT INTEGRATION OVERVIEW
    add_slide("SUBJECT INTEGRATION (6 Domains)", [
        "NLP: TF-IDF and Semantic Reasoning (Llama 3.3).",
        "Web Dev: Flask Backend + Vanilla JS Responsive Dashboard.",
        "Networks: Secure API calls to Groq and DuckDuckGo.",
        "Big Data: Parsing unstructured live news feeds.",
        "DevOps: Modular 4-stage pipeline with fallback systems.",
        "Cybersecurity: Threat scoring and security intelligence logging."
    ])

    # 6. DESIGN AND METHODOLOGY
    add_slide("DESIGN AND METHODOLOGY", [
        "Stage 1 (Local ML): High-speed narrative frame classification.",
        "Stage 2 (NLP): Lexical threat scoring for tone and propaganda.",
        "Stage 3 (News Fetch): Async retrieval of real-world context.",
        "Stage 4 (LLM Reasoning): Llama 3.3 comparative verdict.",
        "Explainability: LIME-based keyword highlighting."
    ], "assets/flowchart.png")

    # 7. SYSTEM ARCHITECTURE
    add_slide("SYSTEM ARCHITECTURE", [
        "Client-Server Model: High-performance dashboard.",
        "Service-Oriented: Independent ML, NLP, and LLM modules.",
        "Data Orchestrator: Python Asyncio for parallel execution.",
        "Persistent Logs: JSON-based security and retraining audit trails."
    ], "assets/architecture.png")

    # 8. IMPLEMENTATION
    add_slide("IMPLEMENTATION", [
        "Backend: Python 3.11 + Flask 3.0.",
        "Frontend: Vanilla HTML5, CSS3, ES6 JavaScript.",
        "ML: Scikit-learn (Logistic Regression).",
        "Cloud AI: Groq SDK for Llama 3.3 70B inference.",
        "Viz: Chart.js for accuracy and latency telemetry."
    ])

    # 9. RESULTS
    add_slide("RESULTS", [
        "Model Accuracy: 83.2% on diverse news benchmarks.",
        "F1-Score: 0.831 | Precision: 0.832.",
        "Inference Latency: 3.4 seconds (Average end-to-end).",
        "Improvement: 85% success in identifying contextual omissions."
    ], "assets/graph.png")

    # 10. INDUSTRY RELEVANCE
    add_slide("INDUSTRY RELEVANCE", [
        "Journalism: Real-time verification of breaking news.",
        "Cybersecurity: Tracking influence operations and narrative attacks.",
        "Strategic Comm: Identifying brand-safety risks and media slant.",
        "Education: Empowering citizens with critical news-reading tools."
    ])

    # 11. CHALLENGES
    add_slide("CHALLENGES", [
        "Linguistic Nuance: Sarcasm and cultural idioms in text.",
        "API Latency: Mitigated via async orchestration.",
        "Search Engine Noise: Filtered via frame-specific queries.",
        "Model Drift: Addressed through an Active Learning feedback loop."
    ])

    # 12. ETHICAL CONSIDERATIONS
    add_slide("ETHICAL CONSIDERATIONS", [
        "Avoiding AI Bias: Ensuring training data diversity.",
        "Privacy: No personal data tracking; focus on public content.",
        "Transparency: Using XAI (LIME) to show 'why' the AI decided.",
        "Accountability: Maintaining audit logs for all analysis verdicts."
    ])

    # 13. FUTURE SCOPE
    add_slide("FUTURE SCOPE", [
        "Multi-modal: Support for Deepfake video and audio detection.",
        "Multi-lingual: Regional language support for local bias.",
        "Blockchain: Immutable fact-checking ledger for trust.",
        "Edge Integration: Browser extensions for real-time protection."
    ])

    # 14. CONCLUSION
    add_slide("CONCLUSION", [
        "The Hybrid AI Dashboard successfully decodes complex media bias.",
        "Combines speed (ML) with depth (LLM) for a robust defense.",
        "Bridges the 'Context Gap' in modern automated fact-checking.",
        "A critical tool for restoring trust in digital information."
    ])

    # 15. REFERENCES
    add_slide("REFERENCES", [
        "GroqCloud (2024): Llama 3.3 Inferencing Reference.",
        "Scikit-learn Documentation: Logistic Regression & TF-IDF.",
        "Jurafsky & Martin (2023): Speech and Language Processing.",
        "Ribeiro et al. (2016): LIME Explainability Framework."
    ])

    filename = "Project_Presentation_Hybrid_AI.pptx"
    prs.save(filename)
    print(f"Presentation generated: {filename}")

if __name__ == "__main__":
    create_presentation()
